# app.py
import os, re, tempfile, fitz, docx, requests, copy
import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR

# ---------------- CONFIG ----------------
GEMINI_API_KEY = "AIzaSyBtah4ZmuiVkSrJABE8wIjiEgunGXAbT3Q"  # 🔑 Hardcoded Gemini API key
TEXT_MODEL_NAME = "gemini-2.0-flash"

# ---------------- GEMINI HELPERS ----------------
def call_gemini(prompt: str) -> str:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{TEXT_MODEL_NAME}:generateContent?key={GEMINI_API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    try:
        resp = requests.post(url, json=payload, timeout=120)
        resp.raise_for_status()
        data = resp.json()
        return data["candidates"][0]["content"]["parts"][0]["text"].strip()
    except Exception as e:
        return f"⚠️ Gemini API error: {e}"

def generate_title(summary: str) -> str:
    prompt = f"""Read the following summary and create a short, clear, presentation-style title.
- Keep it under 10 words
- Do not include birth dates, long sentences, or excessive details
- Just give a clean title, like a presentation heading

Summary:
{summary}
"""
    return call_gemini(prompt).strip()

def extract_slide_count(description: str, default=None):
    m = re.search(r"(\d+)\s*(slides?|sections?|pages?)", description, re.IGNORECASE)
    if m:
        total = int(m.group(1))
        return max(1, total - 1)
    return None if default is None else default - 1

def parse_points(points_text: str):
    points, current_title, current_content = [], None, []
    lines = [re.sub(r"[#*>`]", "", ln).rstrip() for ln in points_text.splitlines()]
    for line in lines:
        if not line or "Would you like" in line: continue
        m = re.match(r"^\s*(Slide|Section)\s*(\d+)\s*:\s*(.+)$", line, re.IGNORECASE)
        if m:
            if current_title:
                points.append({"title": current_title, "description": "\n".join(current_content)})
            current_title, current_content = m.group(3).strip(), []
            continue
        if line.strip().startswith("-"):
            text = line.lstrip("-").strip()
            if text: current_content.append(f"• {text}")
        elif line.strip().startswith(("•", "*")) or line.startswith("  "):
            text = line.lstrip("•*").strip()
            if text: current_content.append(f"- {text}")
        else:
            if line.strip(): current_content.append(line.strip())
    if current_title:
        points.append({"title": current_title, "description": "\n".join(current_content)})
    return points

def generate_outline(description: str):
    num_slides = extract_slide_count(description, default=None)
    if num_slides:
        prompt = f"""Create a PowerPoint outline on: {description}.
Generate exactly {num_slides} content slides (⚠️ excluding the title slide).
Start from Slide 1 as the first *content slide*.
Format:
Slide 1: <Title>
- Bullet
- Bullet
"""
    else:
        prompt = f"""Create a PowerPoint outline on: {description}.
Each slide should have a short title and 3–4 bullet points.
Format:
Slide 1: <Title>
- Bullet
- Bullet
"""
    outline_text = call_gemini(prompt)
    return parse_points(outline_text)

def edit_outline_with_feedback(outline, feedback: str):
    outline_text = "\n".join(
        [f"Slide {i+1}: {s['title']}\n{s['description']}" for i, s in enumerate(outline['slides'])]
    )
    prompt = f"""
    You are an assistant improving a PowerPoint outline.

    Current Outline:
    Title: {outline['title']}
    {outline_text}

    Feedback:
    {feedback}

    Task:
    - Apply the feedback to refine/improve the outline.
    - Return the updated outline with the same format:
      Slide 1: <Title>
      - Bullet
      - Bullet
    - Do NOT add a title slide (I will handle it).
    """
    updated_points = parse_points(call_gemini(prompt))
    return {"title": outline['title'], "slides": updated_points}

def split_text(text: str, chunk_size: int = 8000, overlap: int = 300):
    if not text: return []
    chunks, start, n = [], 0, len(text)
    while start < n:
        end = min(start + chunk_size, n)
        chunks.append(text[start:end])
        if end == n: break
        start = max(0, end - overlap)
    return chunks

# ---------------- UPDATED: Comprehensive summarization ----------------
def summarize_long_text(full_text: str) -> str:
    """
    Produces a comprehensive, exhaustive, and structured summary of the entire document.
    - If the document is short: analyze whole document directly with a single, thorough prompt.
    - If the document is long: split into chunks, produce detailed analysis per chunk, then combine
      into one unified, exhaustive summary that preserves all important points, structure, facts,
      and nuances.
    """
    if not full_text or not full_text.strip():
        return ""

    chunks = split_text(full_text, chunk_size=8000, overlap=400)

    # If it's small, ask Gemini to produce a single exhaustive analysis
    if len(chunks) <= 1:
        prompt = f"""
Read and analyze the entire document below thoroughly. Produce a comprehensive, detailed, and exhaustive summary that preserves every important point, fact, argument, example, and nuance from the text. Do NOT oversimplify or omit material. The output should include:

1) An Executive Summary (one paragraph) that captures the overall purpose and conclusions.
2) A clear reconstruction of the document's structure with headings (e.g., Introduction, Methods/Body, Results/Arguments, Examples, Discussion, Conclusion).
3) For each section: a long, detailed section-by-section summary with important points, supporting evidence, examples, and any arguments or lines of reasoning fully preserved.
4) A consolidated list of Key Facts & Figures (as bullets), including any numbers, dates, named items, or data points.
5) Notable quotes or short excerpts (if present), labelled with approximate location.
6) Any assumptions, limitations, or open questions raised by the document.
7) A final 'Key takeaways' bullet list summarizing the most critical items.

Be exhaustive but keep the final output readable and well-structured. Document:
----------------
{full_text}
----------------
"""
        return call_gemini(prompt).strip()

    # If long, produce detailed analysis for each chunk then combine.
    partial_analyses = []
    for idx, ch in enumerate(chunks, start=1):
        prompt_chunk = f"""
You will be given CHUNK {idx} of a larger document. Carefully analyze this chunk and produce:
A) A detailed, exhaustive summary of CHUNK {idx} that preserves all important points, facts, arguments, examples, and nuance from this chunk.
B) A short heading describing what this chunk contains (e.g., "Introduction", "Methodology", "Case Study", "Analysis", "Conclusion", etc.).
C) A list of Key Facts & Figures found in this chunk (bulleted).
D) Any notable quotes or short excerpts.
E) Any open questions or references that should be cross-referenced with other chunks.

Label the output clearly as "CHUNK {idx} ANALYSIS".

Chunk content follows:
----------------
{ch}
----------------
"""
        analysis = call_gemini(prompt_chunk)
        partial_analyses.append(f"CHUNK {idx} ANALYSIS:\n{analysis.strip()}")

    combined_analyses_text = "\n\n".join(partial_analyses)

    # Combine into one final exhaustive summary
    combine_prompt = f"""
You have a set of detailed chunk analyses from a long document (listed below). Use them to produce ONE unified, coherent, and exhaustive summary of the entire original document. The final output MUST preserve every important point, fact, argument, example, and nuance found across the chunks. DO NOT INVENT new facts.

The final summary should be structured as follows:

1) Executive Summary: One concise paragraph that captures the entire document's purpose and conclusions.
2) Document Structure Reconstruction: Recreate the original document's sections and provide headings (Introduction, Body sections, Results/Arguments, Examples/Case-Studies, Discussion, Conclusion, etc.). For each reconstructed section, provide a thorough, long-form synthesis combining the chunk-level details.
3) Consolidated Key Facts & Figures: A single, deduplicated bulleted list containing all factual items (numbers, dates, names, data points) encountered in the chunks. If a fact appears in multiple chunks, include it once and list chunk locations in parentheses.
4) Important Quotes & Locations: A short list of notable quotes/excerpts and the approximate chunk number where they appear.
5) Assumptions, Limitations, and Open Questions: Combined and organized.
6) Key Takeaways: Clear bulleted summary of the most important conclusions and actionable points.

Below are the chunk analyses. Use them to reconstruct the full document and ensure no detail is lost:

----------------
{combined_analyses_text}
----------------

Now produce the final unified summary described above.
"""
    final_summary = call_gemini(combine_prompt)
    return final_summary.strip()

# ---------------- extract_text and utils ----------------
def extract_text(path: str, filename: str) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        text_parts = []
        doc = fitz.open(path)
        try:
            for page in doc: text_parts.append(page.get_text("text"))
        finally: doc.close()
        return "\n".join(text_parts)
    if name.endswith(".docx"):
        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    if name.endswith(".txt"):
        for enc in ("utf-8","utf-16","utf-16-le","utf-16-be","latin-1"):
            try:
                with open(path,"r",encoding=enc) as f: return f.read()
            except UnicodeDecodeError: continue
        with open(path,"r",encoding="utf-8",errors="ignore") as f: return f.read()
    return ""

def sanitize_filename(name: str) -> str:
    return re.sub(r'[^A-Za-z0-9_.-]', '_', name)

def clean_title_text(title: str) -> str:
    if not title: return "Presentation"
    return re.sub(r"\s+", " ", title.strip())

def hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16))

# ---------------- PPT GENERATOR ----------------
def create_ppt(title, points, filename="output.pptx", title_size=30, text_size=22,
               font="Calibri", title_color="#5E2A84", text_color="#282828", background_color="#FFFFFF"):
    prs = Presentation()
    title = clean_title_text(title)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = hex_to_rgb(background_color)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = textbox.text_frame; tf.word_wrap, tf.auto_size, tf.vertical_anchor = True, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.add_paragraph()
    p.text, p.font.size, p.font.bold, p.font.name, p.font.color.rgb, p.alignment = title, Pt(title_size), True, font, hex_to_rgb(title_color), PP_ALIGN.CENTER

    # Content Slides
    for idx, item in enumerate(points, start=1):
        key_point, description = clean_title_text(item.get("title","")), item.get("description","")
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = hex_to_rgb(background_color)
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1.5))
        tf = textbox.text_frame; tf.word_wrap, tf.auto_size, tf.vertical_anchor = True, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.add_paragraph()
        p.text, p.font.size, p.font.bold, p.font.name, p.font.color.rgb, p.alignment = key_point, Pt(title_size), True, font, hex_to_rgb(title_color), PP_ALIGN.LEFT
        if description:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(5), Inches(4))
            tf = textbox.text_frame; tf.word_wrap = True
            for line in description.split("\n"):
                if line.strip():
                    bullet = tf.add_paragraph()
                    bullet.text, bullet.font.size, bullet.font.name, bullet.font.color.rgb, bullet.level = line.strip(), Pt(text_size), font, hex_to_rgb(text_color), 0
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(8), Inches(0.3))
        tf = textbox.text_frame; p = tf.add_paragraph()
        p.text, p.font.size, p.font.name, p.font.color.rgb, p.alignment = "Generated with AI", Pt(10), font, RGBColor(150,150,150), PP_ALIGN.RIGHT

    prs.save(filename); return filename

# ---------------- STREAMLIT UI ----------------
st.set_page_config(page_title="PPT Generator", layout="wide")
st.title("PPT Generator")

defaults = {
    "messages": [], 
    "outline_chat": None, 
    "summary_text": None, 
    "summary_title": None, 
    "doc_chat_history": [],
    "title_size": 30,
    "text_size": 22,
    "font_choice": "Calibri",
    "title_color": "#5E2A84",
    "text_color": "#282828",
    "bg_color": "#FFFFFF"
}
for k,v in defaults.items():
    if k not in st.session_state: st.session_state[k]=v

# ✅ Always show customization panel
st.subheader("🎨 Customize PPT Style")
col1, col2 = st.columns(2)
with col1: 
    st.session_state.title_size = st.number_input("📌 Title Font Size", 10, 100, st.session_state.title_size)
with col2: 
    st.session_state.text_size = st.number_input("📝 Text Font Size", 8, 60, st.session_state.text_size)

st.session_state.font_choice = st.selectbox(
    "🔤 Font Family",
    ["Calibri","Arial","Times New Roman","Verdana","Georgia","Helvetica","Comic Sans MS"],
    index=["Calibri","Arial","Times New Roman","Verdana","Georgia","Helvetica","Comic Sans MS"].index(st.session_state.font_choice)
)

col3, col4, col5 = st.columns(3)
with col3: 
    st.session_state.title_color = st.color_picker("🎨 Title Color", st.session_state.title_color)
with col4: 
    st.session_state.text_color = st.color_picker("📝 Text Color", st.session_state.text_color)
with col5: 
    st.session_state.bg_color = st.color_picker("🌆 Background Color", st.session_state.bg_color)

# Chat history
for role, content in st.session_state.messages: 
    with st.chat_message(role): st.markdown(content)
for role, content in st.session_state.doc_chat_history:
    with st.chat_message(role): st.markdown(content)

# Upload
uploaded_file = st.file_uploader("📂 Upload a document", type=["pdf","docx","txt"])
if uploaded_file:
    with st.spinner("Processing uploaded file..."):
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(uploaded_file.getvalue()); tmp_path=tmp.name
        text = extract_text(tmp_path, uploaded_file.name); os.remove(tmp_path)
        if text.strip():
            # NOTE: this now produces a full, exhaustive analysis (long-form summary)
            summary = summarize_long_text(text)
            title = generate_title(summary)
            st.session_state.summary_text, st.session_state.summary_title = summary, title
            st.success(f"✅ Uploaded!")
        else:
            st.error("❌ Unsupported, empty, or unreadable file.")

# Chat input
if prompt := st.chat_input("💬 Type a message..."):
    if st.session_state.summary_text:
        if any(w in prompt.lower() for w in ["ppt","slides","presentation"]):
            slides = generate_outline(st.session_state.summary_text + "\n\n" + prompt)
            st.session_state.outline_chat = {"title": st.session_state.summary_title, "slides": slides}
            st.session_state.doc_chat_history.append(("assistant","✅ Generated PPT outline from document."))
        else:
            st.session_state.doc_chat_history.append(("user",prompt))
            reply = call_gemini(f"Answer using only this doc:\n{st.session_state.summary_text}\n\nQ:{prompt}")
            st.session_state.doc_chat_history.append(("assistant",reply))
    else:
        st.session_state.messages.append(("user",prompt))
        if "ppt" in prompt.lower():
            slides = generate_outline(prompt)
            title = generate_title(prompt)
            st.session_state.outline_chat = {"title": title, "slides": slides}
            st.session_state.messages.append(("assistant", f"✅ PPT outline generated! Title: **{title}**"))
        else:
            reply = call_gemini(prompt)
            st.session_state.messages.append(("assistant",reply))
    st.rerun()

# Outline + Feedback + PPT Generation
if st.session_state.outline_chat:
    outline = st.session_state.outline_chat
    st.subheader(f"📝 Preview Outline: {outline['title']}")
    for idx,slide in enumerate(outline["slides"],start=1):
        with st.expander(f"Slide {idx}: {slide['title']}",expanded=False):
            st.markdown(slide["description"].replace("\n","\n\n"))

    new_title = st.text_input("📌 Edit Title", value=outline.get("title","Untitled"))
    feedback_box = st.text_area("✏️ Feedback for outline (optional):")

    col6, col7 = st.columns(2)

    with col6:
        if st.button("🔄 Apply Feedback"):
            with st.spinner("Updating outline with feedback..."):
                try:
                    updated_outline = edit_outline_with_feedback(outline, feedback_box)
                    updated_outline["title"] = new_title.strip() if new_title else updated_outline["title"]
                    st.session_state.outline_chat = updated_outline
                    st.success("✅ Outline updated with feedback!")
                    st.rerun()
                except Exception as e:
                    st.error(f"❌ Feedback error: {e}")

    with col7:
        if st.button("✅ Generate PPT"):
            with st.spinner("Generating PPT..."):
                filename = f"{sanitize_filename(new_title)}.pptx"
                create_ppt(new_title, outline["slides"], filename,
                           title_size=int(st.session_state.title_size), 
                           text_size=int(st.session_state.text_size),
                           font=st.session_state.font_choice, 
                           title_color=st.session_state.title_color,
                           text_color=st.session_state.text_color, 
                           background_color=st.session_state.bg_color)
                with open(filename,"rb") as f:
                    st.download_button("⬇️ Download PPT", data=f, file_name=filename,
                                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
